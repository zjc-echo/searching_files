import os
import fitz
from docx import Document
from pptx import Presentation
from docxtpl import DocxTemplate
import win32com.client
#import easyocr
import time

class Files:
    def __init__(self, folder_path):
        self.folder_path = folder_path
        self.file_list = []
        self.file_name = ""

    def get_file_extension(self, file_name):
        extension = os.path.splitext(file_name)[1].lower()
        return extension

    def get_file_details(self):
        pass

    def get_file_address(self):
        file_address = self.folder_path+'\\'+self.file_name
        return file_address

class FileDOCX(Files):
    def get_file_details(self):
        # 调用父类的方法，获取完整的文件路径
        file_path = self.get_file_address()
        # 将 Word 文档转换为 PDF
        pdf_path = self.convert_to_pdf(file_path)
        # 使用 fitz 模块打开 PDF 文件
        doc = fitz.open(pdf_path)
        # 遍历每一页，获取页码和内容
        file_details = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text('text').strip()
            file_details.append([page_num, page_text])
        # 关闭 PDF 文件
        doc.close()
        # 删除生成的 PDF 文件
        os.remove(pdf_path)
        # 将页码和内容存储在 self.file_list 中
        self.file_list = file_details
        return

    def convert_to_pdf(self, file_path):
        # 创建 Word 应用程序对象并打开文档
        app_names = ["kwps.Application", "Word.Application", "wps.Application"]
        word_app = None
        for app_name in app_names:
            try:
                word_app = win32com.client.Dispatch(app_name)
                break
            except Exception:
                pass

        if word_app is None:
            raise Exception("No compatible Word application found")
        doc = word_app.Documents.Open(file_path)
        # 将文档保存为 PDF
        pdf_path = file_path.replace(".docx", ".pdf")
        doc.SaveAs(pdf_path, FileFormat=17)
        # 关闭 Word 文档和应用程序
        doc.Close()
        word_app.Quit()
        return pdf_path

class FilePPTX(Files):
    def get_file_details(self):
        # 调用父类的方法，获取完整的文件路径
        file_path = self.get_file_address()

        # 使用 pptx 模块打开 PPTX 文件
        presentation = Presentation(file_path)

        # 遍历每一页，获取页码和内容
        file_details = []
        for slide_num, slide in enumerate(presentation.slides):
            slide_content = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_content += paragraph.text.strip() + "\n"
            file_details.append([slide_num, slide_content])

        # 将页码和内容存储在 self.file_list 中
        self.file_list = file_details

        # 关闭 PPTX 文件
        presentation = None

class FilePDF(Files):
    def get_file_details(self):
        # 调用父类的方法，获取完整的文件路径
        file_path = self.get_file_address()
        # 使用 fitz 模块打开 PDF 文件
        doc = fitz.open(file_path)
        # 遍历每一页，获取页码和内容
        file_details = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text('text').strip()
            file_details.append([page_num, page_text])
        self.file_list = file_details
        # 关闭 PDF 文件
        doc.close()
        # 将页码和内容存储在 self.file_list 中
        return


class FilePDF_ReadUnable(FilePDF):
    '''
    def get_file_details(self):
        file_path = self.get_file_address()
        doc = fitz.open(file_path)
        file_details = []
        # 创建 EasyOCR 实例，设置语言参数为中英双语
        reader = easyocr.Reader(['en', 'ch_sim'], gpu=True)
        # 遍历每一页，将每一页转换为图片并进行文字识别
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            # 将每一页转换为图片
            pix = page.get_pixmap()
            # 保存图片到临时文件
            image_path = f"page_{page_num}.png"
            pix.save(image_path)
            # 进行文字识别
            result = reader.readtext(image_path, detail=0)
            # 删除临时图片文件
            os.remove(image_path)
            # 将页码和识别的文字内容添加到文件详情列表中
            file_details.append([page_num, result])
        doc.close()
        self.file_list = file_details
        return
    '''

def main():
    pass

if __name__ == "__main__":
    main()